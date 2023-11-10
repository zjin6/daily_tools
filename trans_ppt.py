import pptx
from googletrans import Translator

# Open the PowerPoint file
ppt = pptx.Presentation(r"C:\Users\zjin6\Downloads\锂电池讲解材料-Anna.pptx")

# Create a Translator object
translator = Translator()



# Save the translated PowerPoint file
ppt.save(r'C:\Users\zjin6\Downloads\锂电池translated.pptx')

# Loop through each slide in the PowerPoint file
for slide in ppt.slides:
    # Loop through each shape on the slide
    for shape in slide.shapes:
        # If the shape has text, translate it to English
        if shape.has_text_frame:
            # Get the original text
            text = shape.text_frame.text
            # Translate the text to English
            try:
                translation = translator.translate(text, dest='en').text
            except:
                # Handle the error here (e.g., print an error message)
                print('Translation error')
                # Set the translation to the original text
                translation = text
            # Replace the text in the shape with the translated text
            shape.text_frame.text = translation

ppt.save(r'C:\Users\zjin6\Downloads\锂电池translated.pptx')