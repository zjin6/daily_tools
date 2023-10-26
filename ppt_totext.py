import pptx

# Open the PowerPoint file
ppt = pptx.Presentation(r"C:\Users\zjin6\Downloads\锂电池讲解材料-Anna.pptx")

# Create an empty list to store the text from each slide
text_list = []

# Loop through each slide in the PowerPoint file
for slide in ppt.slides:
    # Create an empty string to store the text from each shape on the slide
    slide_text = ''
    # Loop through each shape on the slide
    for shape in slide.shapes:
        # If the shape has text, add it to the slide_text string
        if shape.has_text_frame:
            slide_text += shape.text_frame.text
    # Add the slide_text string to the text_list
    text_list.append(slide_text)

# Print the text_list
for i, text in enumerate(text_list):
    print('page' + str(i))
    print(text)
